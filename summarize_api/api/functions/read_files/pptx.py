from pptx import Presentation
from utils.logs.filehandler import logging
import os
import tempfile

async def read_pptx(file_content: bytes) -> str:
    try:
        output_dir = os.path.join(os.getcwd(), "utils", "output")
        output_text_file = os.path.join(output_dir, "texto.txt")

        logging.info("Extrayendo texto de la presentación PowerPoint")
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
            temp_file.write(file_content)
            temp_file_path = temp_file.name

        text = extract_text_from_presentation(temp_file_path)
        save_text_to_file(text, output_text_file)

        logging.debug(f"Texto extraído del archivo y guardado en: {output_text_file}")
        os.unlink(temp_file_path)
        
        return text
    except Exception as e:
        logging.error(f"Error al extraer texto del archivo PowerPoint: {e}", exc_info=True)
        return str(e)

def create_output_directory() -> str:
    output_dir = os.path.join(os.getcwd(), "utils", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    return output_dir

def extract_text_from_presentation(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""

    for slide_number, slide in enumerate(prs.slides):
        text += f"\nSlide {slide_number + 1}:\n"
        
        for shape in slide.shapes:
            # Extraer texto de cuadros de texto y figuras
            if hasattr(shape, "text") and shape.has_text_frame:
                text += extract_text_from_text_frame(shape.text_frame)
            # Extraer texto de tablas
            elif shape.has_table:
                text += extract_text_from_table(shape.table)

    return text

def extract_text_from_text_frame(text_frame) -> str:
    text = ""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            text += run.text + "\n"
    return text

def extract_text_from_table(table) -> str:
    text = ""
    for row in table.rows:
        for cell in row.cells:
            text += cell.text + " "
        text += "\n"
    return text

def save_text_to_file(text: str, file_path: str):
    try:
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(text)
    except Exception as e:
        logging.error(f"Error al guardar el texto en el archivo: {e}", exc_info=True)
        raise